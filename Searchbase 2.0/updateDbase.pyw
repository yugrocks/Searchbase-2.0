import os,time
import pickle
from lxml import _elementpath as _dummy
from pptx import Presentation
import openpyxl
import docx
import re
import details_song

found=False
words_to_be_deleted=[]
most_recent=0
most_recent_current=0
y=[]
list2=None
xdct={}
listofdocs=[]
listoffiles=[]
metaclass=None
sentinel=False

def add_into_list(word):
    global list2
    if not word in list2:
       for _ in list2:
          if _ > word:
              list2.insert(list2.index(_),word)
              break


def retrieve_file_dbase():
    global xdct
    if not os.path.exists(r'C:\database2'):
        os.mkdir(r'C:\database2')
        return
    os.chdir(r'C:\database2')
    for _ in range(1,11):
            fileopen=open(r'C:\database2\dbw_{}.pkl'.format(_),'rb')
            xdct.update(pickle.load(fileopen))
            fileopen.close()


def getmeta(file):
   global metaclass
   print(metaclass)
   if metaclass is not None:
       meta=metaclass.getArtist(file).lower()+" "+metaclass.getGenre(file).lower()
       return remove(meta)
   return ""


b_ref=[',','<','.','>','/','\',','?','"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]
def remove(s):
    s=s.lower()
    list1=[]
    k=""
    for _ in s:
        if _ is not "'":
          if _ in b_ref:
            if k!="":
                list1.extend(k.replace("'","").split())
            k=""
            continue
        k+=_
    list1.extend(k.replace("'","").split())
    return list1


def database(dir):
    global most_recent_current
    b=os.listdir(dir)
    for entity in b:
        if sentinel:
            break
        time.sleep(0.009)
        if entity=="$Recycle.Bin" or entity=="$RECYCLE.BIN" or entity=="System Volume Information" or "appdata" in entity.lower() or entity.lower()=='windows':
           continue
        if os.path.isdir(r'{}\{}'.format(dir,entity)):
            try:
                if os.path.getmtime(r'{}\{}'.format(dir,entity)) >most_recent:
                    y.append(os.path.getmtime(r'{}\{}'.format(dir,entity)))
                    processFile(os.path.normpath(r'{}\{}'.format(dir,entity)))
                database(r'{}\{}'.format(dir,entity))
            except:
                continue
        else:
            try:
                if os.path.getmtime(os.path.join(dir,entity))>most_recent:
                    y.append(os.path.getmtime(r'{}\{}'.format(dir,entity)))
                    processFile(os.path.normpath(r'{}\{}'.format(dir,entity)))
            except:
                continue


def get_ext(s):
    s=s.split(".")
    if len(s)==1:
       return ""
    return s[len(s)-1]


def check_directory(dir):
    global found,x,sentinel
    for file in os.listdir(dir):
        if sentinel:
            break
        time.sleep(0.01)
        if os.path.isdir(os.path.join(dir,file)):
            ext='folder'
        else:
            ext=get_ext(file)
        if ext=='mp3':
                print("I am here")
                meta=getmeta(os.path.join(dir,file))
                for artist in meta:
                    print(artist)
                    if artist.lower() in x:
                       found=False
                       for tpl in x[artist.lower()]:
                           time.sleep(0.001)
                           if os.path.normpath(os.path.join(dir,file)) == listoffiles[tpl[0]]:
                               found=True
                       if not found:
                         if not (os.path.normpath(os.path.join(dir,file)),ext) in listoffiles:
                           listoffiles.append(os.path.normpath(os.path.join(dir,file)))
                         x[artist.lower()].append((listoffiles.index(os.path.normpath(os.path.join(dir,file))),ext))
                         add_into_list(artist.lower())
                         print("artist successfully added to database {}".format(artist))
                         found=False
                    else:
                        if not (os.path.normpath(os.path.join(dir,file)),ext) in listoffiles:
                                listoffiles.append(os.path.normpath(os.path.join(dir,file)))
                        listn=[];
                        listn.append((listoffiles.index(os.path.normpath(os.path.join(dir,file))),ext))
                        x[artist.lower()]=listn
                        add_into_list(artist.lower())
                        print("File successfully added to database{}".format(file))
        for word in remove(file):
            if word.lower() in x:
                found=False
                for tpl in x[word.lower()]:
                    time.sleep(0.0001)
                    if os.path.normpath(os.path.join(dir,file)) == listoffiles[tpl[0]]:
                        found=True
                if not found:
                    if not (os.path.normpath(os.path.join(dir,file)),ext) in listoffiles:
                        listoffiles.append(os.path.normpath(os.path.join(dir,file)))
                    x[word.lower()].append((listoffiles.index(os.path.normpath(os.path.join(dir,file))),ext))
                    add_into_list(word.lower())
                    #print("File successfully added to database {}".format(file))
                    found=False
            else:
                if not (os.path.normpath(os.path.join(dir,file)),ext) in listoffiles:
                        listoffiles.append(os.path.normpath(os.path.join(dir,file)))
                listn=[];
                listn.append((listoffiles.index(os.path.normpath(os.path.join(dir,file))),ext))
                x[word.lower()]=listn
                add_into_list(word.lower())
                #print("File successfully added to database{}".format(file))


def check_file(file):
    print("checking file")
    scan2(file)


def processFile(file):
    global most_recent
    if os.path.isdir(file):
        print("directory {} has been modified".format(os.path.basename(file)))
        check_directory(file)
    else:
        check_file(file)
    jingalala2=open(r'C:\database2\DATABASER.pkl','wb')
    pickle.dump({"most_recent":most_recent},jingalala2)
    jingalala2.close()

def scan2(entity):
     global xdct,listofdocs
     scannable=False
     entity=os.path.normpath(entity)
     file_contents=[]
     print("scanning : ",entity)
     stopwordlist=['None','is','an','the','are','a','of','and','to','for','in','it','',' ']
     if '.docx' in os.path.basename(entity):
        scannable=True
        time.sleep(0.01)
        if not os.path.normpath(entity) in listofdocs:
            listofdocs.append(os.path.normpath(entity))
        try:
           a=docx.Document(entity)
           for line in a.paragraphs:
              time.sleep(0.005)
              b=remove(line.text)
              for word in b:
                 time.sleep(0.005)
                 if word=='' or word==' ':
                    continue
                 elif '\u2019' in word:
                    wordn=word.replace('\u2019','')
                 elif '\u2026' in word:
                    wordn=word.replace('\u2026','')
                 else:
                    wordn=word.lower()
                 file_contents.append(wordn)
        except:
           print('no')
     elif '.txt' in os.path.basename(entity) or '.cpp' in os.path.basename(entity) or '.csv' in os.path.basename(entity):
        scannable=True
        time.sleep(0.01)
        if not os.path.normpath(entity) in listofdocs:
            listofdocs.append(os.path.normpath(entity))
        try:
           am=open(entity)
           for word in am:
               time.sleep(0.005)
               list2=remove(word)
               for e in list2:
                  d=e.replace('\x00','').lower()
                  if d=='' or d==' ' or d in stopwordlist:
                     continue
                  file_contents.append(d)
        except:
           print('no')
     elif '.pptx' in os.path.basename(entity):
        scannable=True
        time.sleep(0.001)
        if not os.path.normpath(entity) in listofdocs:
            listofdocs.append(os.path.normpath(entity))
        try:
            prs = Presentation(entity)
            for slide in prs.slides:
               time.sleep(0.005)
               for shape in slide.shapes:
                  if not shape.has_text_frame:
                     continue
                  for paragraph in shape.text_frame.paragraphs:
                     time.sleep(0.005)
                     a=remove(paragraph.text)
                     for ptword in a:
                         time.sleep(0.005)
                         if ptword in stopwordlist:
                           continue
                         file_contents.append(ptword.lower())
        except:
           print('no')
     elif '.xlsx' in os.path.basename(entity):
        scannable=True
        time.sleep(0.001)
        if not os.path.normpath(entity) in listofdocs:
            listofdocs.append(os.path.normpath(entity))
        try:
           wb=openpyxl.load_workbook(entity)
           getnm=wb.get_sheet_names()
           for sheetn in getnm:
              time.sleep(0.005)
              sheet=wb.get_sheet_by_name(sheetn)
              hr=sheet.max_row
              hc=sheet.max_column
              for i in range(1,hr):
                 time.sleep(0.005)
                 for j in range(1,hc):
                    clvalue2=sheet.cell(row=i,column=j).value
                    clvalue="{}".format(clvalue2)
                    if type(clvalue) is not int:
                       if clvalue in stopwordlist or clvalue is None or '1' in clvalue or '2' in clvalue or '3' in clvalue or '4' in clvalue or '5' in clvalue or '6' in clvalue or '7' in clvalue or '8' in clvalue or '9' in clvalue or '0' in clvalue:
                           continue
                    else:
                       continue
                    if type(clvalue) is str:
                       clvalue2=remove(clvalue)
                       for ant in clvalue2:
                           file_contents.append(ant.lower())
        except:
            return
     if scannable:
         retrieve_file_dbase()
         for _ in xdct:
             time.sleep(0.001)
             if not _ in file_contents:
                 if listofdocs.index(os.path.normpath(entity)) in xdct[_]:
                     xdct[_].remove(listofdocs.index(os.path.normpath(entity)))
         for _ in file_contents:
                 if _ in xdct:
                     if listofdocs.index(os.path.normpath(entity)) in xdct[_]:
                         continue
                 else:
                     xdct[_]=[listofdocs.index(os.path.normpath(entity))]
         print('Done')
         store_dbase2(xdct)
     jingalala2=open(r'C:\database2\DATABASER.pkl','wb')
     pickle.dump({"most_recent":most_recent},jingalala2)
     jingalala2.close()

def start(dct,recent,liste,listofdocs2,listoffiles2,metaclas):
  global x,most_recent,most_recent_current,y,list2,listofdocs,listoffiles,sentinel
  global metaclass
  metaclass=metaclas
  listofdocs=listofdocs2
  listoffiles=listoffiles2
  x=dct
  list2=liste
  most_recent=recent
  most_recent_current=most_recent
  while True:
      if sentinel:
            break
      print(len(x))
      time.sleep(0.01)
      ll2=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)
      for direc in ll2:
          try:
              ll3=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk where drivetype=2 get deviceid").read(),re.MULTILINE)
              if direc in ll3:
                  continue
              rcheck=os.popen("vol {}".format(direc))
              rchk=rcheck.readline()
              if 'recovery' in rchk or 'RECOVERY' in rchk:
                  continue
              try:
                 database(direc+"\\")
              except:
                  continue
          except:
              continue
      if len(y)>0:
          most_recent=max(y)
      y=[]
      check_for_deletions(x,list2,listoffiles)
      store_dbase(x)
      jingalala2=open(r'C:\database2\DATABASER.pkl','wb')
      pickle.dump({"most_recent":most_recent},jingalala2)
      jingalala2.close()


def check_for_deletions(x,list2,listoffiles):
    global sentinel
    for _ in x:
        if sentinel:
            break
        time.sleep(0.01)
        for tpl in x[_]:
            time.sleep(0.005)
            if not os.path.exists(listoffiles[tpl[0]]):
                for word in remove(os.path.basename(listoffiles[tpl[0]])):
                    word=word.lower()
                    if word in x:
                        try:
                            x[word].remove((tpl[0],get_ext(listoffiles[tpl[0]])))
                        except:
                            try:
                               x[word].remove((tpl[0],'folder'))
                            except:
                                print("")
                        print(listoffiles[tpl[0]], ' removed from ',word)
                        if len(x[word])==0:
                            words_to_be_deleted.append(word)
                            print(word," deleted")
                listoffiles[tpl[0]]=""
    for word in words_to_be_deleted:
        if word in x and word in list2:
            list2.remove(word)
            x.pop(word)

def store_dbase(x,parts=10):
    if not os.path.exists(r'C:\database2'):
        os.mkdir(r'C:\database2')
    os.chdir(r'C:\database2')
    av_length=int(len(x)/parts)
    counter=0
    rounds=0
    dct_int={}
    for k in x:
        counter+=1
        dct_int[k]=x[k]
        if counter==av_length and rounds<parts-1:
            counter=0
            rounds+=1
            fileopen=open('dbf_{}.pkl'.format(rounds),'wb')
            pickle.dump(dct_int,fileopen)
            fileopen.close()
            dct_int.clear()
    fileopen=open('dbf_{}.pkl'.format(parts),'wb')
    pickle.dump(dct_int,fileopen)
    fileopen.close()
    dct_int.clear()
    fileopen2=open('dblst2.pkl','wb')
    pickle.dump(listoffiles,fileopen2)
    fileopen2.close()
    print('database has been saved')


def store_dbase2(xdct, parts=10):
    if not os.path.exists(r'C:\database2'):
        os.mkdir(r'C:\database2')
    os.chdir(r'C:\database2')
    av_length=int(len(xdct)/parts)
    counter=0
    rounds=0
    dct_int={}
    for k in xdct:
        counter+=1
        dct_int[k]=xdct[k]
        if counter==av_length and rounds<parts-1:
            counter=0
            rounds+=1
            fileopen2=open('dbw_{}.pkl'.format(rounds),'wb')
            pickle.dump(dct_int,fileopen2)
            fileopen2.close()
            dct_int.clear()
    fileopen2=open('dbw_{}.pkl'.format(parts),'wb')
    pickle.dump(dct_int,fileopen2)
    fileopen2.close()
    dct_int.clear()
    xdct.clear()
    fileopen2=open('dblst.pkl','wb')
    pickle.dump(listofdocs,fileopen2)
    fileopen2.close()
    jingalala2=open(r'C:\database2\DATABASER.pkl','wb')
    pickle.dump({"most_recent":most_recent},jingalala2)
    jingalala2.close()

