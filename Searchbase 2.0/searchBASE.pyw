import os
import time
import operator
from tkinter import *
import threading
import searchbasenew
import spellCheck2
from functools import partial
from lxml import _elementpath as _dummy
from pptx import Presentation
import openpyxl
import pickle
import docx
import updateScreen
import shortcut
import updateDbase
import settings
import details_song

supremedir=None
metaclass=None
iterations=0
alt_count=0
a=None
x={}
xcount={}
xdct={}
list2=[]
alt_UI=False
alt_Frame=None
alt_Canvas=None
frame1=None
frame2=None
frme=None
frme2=None
frame3=None
frame4 =None
lm=None
ld=None
lf=None
most_recent=0
updateScrn=None
textLabel=None
thisdict2={}
current_query=""
memory={}
shrt=None
shortcut_active=False
listofdocs=[]
listoffiles=[]
helper_dict={}

b_ref=[',','<','.','>','/','\',','?','"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]

def capspermutation(string):
   c=[]
   l=[]
   st3=string.upper()
   st2=string.lower()
   c.append(st2)
   c.append(st3)
   c.append(string.lower().replace(string[0],string[0].upper(),1))
   return(c)


def remove(s):
    list1=[]
    k=""
    for _ in s:
        if _ is not "'":
          if _ in b_ref:
            if k!="":
                list1.extend(k.replace("'","").split())
            k=""
            continue
        k+=_
    list1.extend(k.replace("'","").split())
    return list1

def get_ext(s):
    s=s.split(".")
    if len(s)==1:
       return ""
    return s[len(s)-1]


def initgetmeta():
   global metaclass
   try:
      metaclass=details_song.getMeta()
   except:
      metaclass=None


def getmeta(file):
   global metaclass
   if metaclass is not None:
       meta=metaclass.getArtist(file).lower()+" "+metaclass.getGenre(file).lower()
       return remove(meta)
   return ""

#DATABASE CREATION STARTS HERE_______________________________________________________

index=-1

def database(dir):
    global most_recent,updateScrn,listoffiles,index,helper_dict
    b=os.listdir(dir)
    try:
       textLabel.config(text="Scanning\n{}".format(dir))
    except:
       print('')
    if os.path.getmtime(dir)>most_recent:
        most_recent=os.path.getmtime(dir)
    for entity in b:
        if entity=="$Recycle.Bin" or entity=="$RECYCLE.BIN" or entity=="System Volume Information" or "appdata" in entity.lower() or entity.lower()=='windows':#some useless directories to look into
            continue
        if not os.path.normpath(os.path.join(dir,entity)) in helper_dict:
            helper_dict[os.path.normpath(os.path.join(dir,entity))]=1
            listoffiles.append(os.path.normpath(os.path.join(dir,entity)))
            index+=1
        try:
            if os.path.getmtime(r'{}\{}'.format(dir,entity))>most_recent:
                most_recent=os.path.getmtime(dir)
        except:
            print("access denied")
        if os.path.isdir(r'{}\{}'.format(dir,entity)):
            for _ in remove(entity):
                if _.lower() in x:
                    if not (index,"folder") in x[_.lower()]:
                        x[_.lower()].append((index,"folder"))
                    continue
                listn=[];listn.append((index,"folder"))
                x[_.lower()]=listn
            try:
                database(r'{}\{}'.format(dir,entity))
            except:
                continue
        else:
            ext=get_ext(entity)
            if ext=='mp3':
              try:
                for artist in getmeta(r'{}\{}'.format(dir,entity)):
                    if artist.lower() in x:
                        if not (index,ext) in x[artist.lower()]:
                            x[artist.lower()].append((index,ext))
                    else:
                        x[artist.lower()]=[(index,ext)]
              except:
                  print("no meta available")
            for _ in remove(entity):
                if _.lower() in x:
                    if not (index,ext) in x[_.lower()]:
                        x[_.lower()].append((index,ext))
                    continue
                listn=[];listn.append((index,ext))
                x[_.lower()]=listn


def updateDB():
    global updateScrn,textLabel,supremedir
    initgetmeta()
    updateScrn=updateScreen.UpdateScreen(supremedir)
    textLabel=updateScrn.getLabel()
    threading.Thread(target=create_dbase).start()
    updateScrn.root2.mainloop()

def create_dbase():
      global updateScrn,helper_dict
      ll2=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk get deviceid").read(),re.MULTILINE)
      for direc in ll2:
          try:
              ll3=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk where drivetype=2 get deviceid").read(),re.MULTILINE)
              if direc in ll3:
                  continue
              rcheck=os.popen("vol {}".format(direc))
              rchk=rcheck.readline()
              if 'recovery' in rchk or 'RECOVERY' in rchk:
                  continue
              try:
                 print(direc+"\\")
                 database(direc+"\\")
              except:
                  continue
          except:
              continue
      for direc in ll2:
          try:
              ll3=re.findall(r"[A-Z]+:",os.popen("wmic logicaldisk where drivetype=2 get deviceid").read(),re.MULTILINE)
              ll3+=[os.path.splitdrive(os.path.expanduser("~"))[0]]
              if direc in ll3:
                  continue
              rcheck=os.popen("vol {}".format(direc))
              rchk=rcheck.readline()
              if 'recovery' in rchk or 'RECOVERY' in rchk:
                  continue
              try:
                 scan2(direc)
              except:
                  continue
          except:
              continue
      scan2(os.path.expanduser("~"))
      helper_dict={}
      helper_dict=None
      store_dbase(10)
      updateScrn.label1.config(text="The database has been created.\nYou should close this and restart the app.")
      textLabel.config(text="")


def store_dbase(parts=7):
    global most_recent,list2,x,listoffiles
    if not os.path.exists(r'C:\database2'):
        os.mkdir(r'C:\database2')
    os.chdir(r'C:\database2')
    av_length=int(len(x)/parts)
    counter=0
    rounds=0
    dct_int={}
    for k in x:
        counter+=1
        dct_int[k]=x[k]
        if counter==av_length and rounds<parts-1:
            counter=0
            rounds+=1
            fileopen=open('dbf_{}.pkl'.format(rounds),'wb')
            pickle.dump(dct_int,fileopen)
            fileopen.close()
            dct_int.clear()
    fileopen=open('dbf_{}.pkl'.format(parts),'wb')
    pickle.dump(dct_int,fileopen)
    fileopen.close()
    dct_int.clear()
    av_length=int(len(xdct)/parts)
    counter=0
    rounds=0
    dct_int={}
    for k in xdct:
        counter+=1
        dct_int[k]=xdct[k]
        if counter==av_length and rounds<parts-1:
            counter=0
            rounds+=1
            fileopen2=open('dbw_{}.pkl'.format(rounds),'wb')
            pickle.dump(dct_int,fileopen2)
            fileopen2.close()
            dct_int.clear()
    fileopen2=open('dbw_{}.pkl'.format(parts),'wb')
    pickle.dump(dct_int,fileopen2)
    fileopen2.close()
    dct_int.clear()
    fileopen2=open('dblst.pkl','wb')
    pickle.dump(listofdocs,fileopen2)
    fileopen2.close()
    fileopen2=open('dblst2.pkl','wb')
    pickle.dump(listoffiles,fileopen2)
    fileopen2.close()
    file_for_recent=open('DATABASER.pkl','wb')
    pickle.dump({"most_recent":most_recent},file_for_recent)
    file_for_recent.close()
    list2=list(x.keys())
    list2.sort()


def retrieve_dbase():
    global x,list2,most_recent,thisdict2,listofdocs,listoffiles
    os.chdir(r'C:\database2')
    iterations=0
    try:
        for _ in range(1,11):
            fileopen=open(r'C:\database2\dbf_{}.pkl'.format(_),'rb')
            x.update(pickle.load(fileopen))
            fileopen.close()
        list2=list(x.keys())
        list2.sort()
        fileopen2=open(r'C:\database2\dblst2.pkl','rb')
        listoffiles=pickle.load(fileopen2)
        fileopen2.close()
        fileopen2=open(r'C:\database2\DATABASER.pkl','rb')
        most_recent0=pickle.load(fileopen2)
        most_recent=most_recent0["most_recent"]
        fileopen2.close()
        for _ in range(1,11):
            fileopen=open(r'C:\database2\dbw_{}.pkl'.format(_),'rb')
            thisdict2.update(pickle.load(fileopen))
            fileopen.close()
        fileopen2=open(r'C:\database2\dblst.pkl','rb')
        listofdocs=pickle.load(fileopen2)
        fileopen2.close()
    except:
        print('no')
        updateDB()


stopwordlist=['none','is','an','the','are','a','of','and','to','for','in','it','',' ']
def scan2(root_dir):
    files=os.listdir(root_dir)
    for file in files:
        if file=="$Recycle.Bin" or file=="$RECYCLE.BIN" or file=="System Volume Information" or "appdata" in file.lower():
           continue
        if not os.path.isdir(os.path.join(root_dir,file)):
            if get_ext(file)=='docx':
                if not os.path.normpath(os.path.join(root_dir,file)) in listofdocs:
                    listofdocs.append(os.path.normpath(os.path.join(root_dir,file)))
                try:
                     textLabel.config(text="Scanning\n{}".format(file))
                except:
                    print('')
                try:
                    a=docx.Document(r'{}\\{}'.format(root_dir,file))
                    for line in a.paragraphs:
                        b=remove(line.text)
                        for word in b:
                            if word=='' or word==' ':
                                continue
                            elif '\u2019' in word:
                                wordn=word.replace('\u2019','').lower()
                            elif '\u2026' in word:
                                wordn=word.replace('\u2026','').lower()
                            else:
                                wordn=word.lower()
                            if wordn in xdct:
                                if listofdocs.index(os.path.normpath(os.path.join(root_dir,file))) in xdct[wordn]:
                                    continue
                                xdct[wordn].append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                print(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                continue
                            listm=[]
                            listm.append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                            xdct[wordn]=listm
                except:
                    continue
            elif get_ext(file)=='txt' or get_ext(file)=='cpp' or get_ext(file)=='csv':
                if not os.path.normpath(os.path.join(root_dir,file)) in listofdocs:
                    listofdocs.append(os.path.normpath(os.path.join(root_dir,file)))
                try:
                     textLabel.config(text="Scanning\n{}\\{}".format(root_dir,file))
                except:
                    print('')
                try:
                    am=open(r'{}\\{}'.format(root_dir,file))
                    for word in am:
                        list2=remove(word)
                        for e in list2:
                            d=e.replace('\x00','').lower()
                            if d=='' or d==' ' or d in stopwordlist:
                                continue
                            if d in xdct:
                                if listofdocs.index(os.path.normpath(os.path.join(root_dir,file))) in xdct[d]:
                                    continue
                                xdct[d].append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                continue
                            listn=[]
                            listn.append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                            xdct[d]=listn
                except:
                    continue
            elif get_ext(file)=='pptx':
                if not os.path.normpath(os.path.join(root_dir,file)) in listofdocs:
                    listofdocs.append(os.path.normpath(os.path.join(root_dir,file)))
                try:
                     textLabel.config(text="Scanning\n{}".format(file))
                except:
                    print('')
                try:
                    prs = Presentation(r'{}\\{}'.format(root_dir,file))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                a=remove(paragraph.text)
                                for ptword in a:
                                    if ptword.lower() in stopwordlist:
                                        continue
                                    if ptword.lower() in xdct:
                                        if listofdocs.index(os.path.normpath(os.path.join(root_dir,file))) in xdct[ptword]:
                                            continue
                                        xdct[ptword.lower()].append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                        continue
                                    listn=[]
                                    listn.append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                    xdct[ptword.lower()]=listn
                except:
                    continue
            elif get_ext(file)=='xlsx':
                if not os.path.normpath(os.path.join(root_dir,file)) in listofdocs:
                    listofdocs.append(os.path.normpath(os.path.join(root_dir,file)))
                try:
                     textLabel.config(text="Scanning\n{}".format(file))
                except:
                    print('')
                try:
                    wb=openpyxl.load_workbook(r'{}\\{}'.format(root_dir,file))
                    getnm=wb.get_sheet_names()
                    for sheetn in getnm:
                        sheet=wb.get_sheet_by_name(sheetn)
                        hr=sheet.max_row
                        hc=sheet.max_column
                        for i in range(1,hr):
                            for j in range(1,hc):
                                clvalue2=sheet.cell(row=i,column=j).value
                                clvalue="{}".format(clvalue2)
                                if type(clvalue) is not int:
                                    if clvalue.lower() in stopwordlist or clvalue is None:
                                        continue
                                else:
                                    continue
                                if type(clvalue) is str:
                                    clvalue2=remove(clvalue)
                                    for ant in clvalue2:
                                        if ant.lower() in xdct:
                                            xdct[ant.lower()].append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                            continue
                                        listn=[]
                                        listn.append(listofdocs.index(os.path.normpath(os.path.join(root_dir,file))))
                                        xdct[ant.lower()]=listn
                except:
                    continue
        else:
          try:
            scan2(r'{}\\{}'.format(root_dir,file))
          except:
            continue


#DATABASE FUNCTIONS END HERE________________________________________________________________________

#================================================================================

#SEARCH FUNCTIONS START HERE____________________________________________________________


def _on_mousewheel(frame,event):
       try:
          frame.yview_scroll(-1*int(event.delta/120), "units")
       except:
          pass


listofwords=[]
points_dict={}
stopsign=False


def process_indx():
   global frame1,frame2,frme,frme2,frame3,frame4,memory,current_query,listoffiles
   folder_count=0;media_count=0;doc_count=0
   for _ in listofwords:
      for ent in x[list2[_]]:
         if ent in points_dict:
            points_dict[ent]+=1
            continue
         points_dict[ent]=1
   for _ in current_query.split():
       if _ in memory:
           for file in memory[_]:
              if len(file)>1:
                  if 'int' in str(type(file[0])):
                      if (file[0],get_ext(listoffiles[file[0]])) in points_dict:
                          points_dict[(file[0],get_ext(listoffiles[file[0]]))]+=1*file[1]
                      if (file[0],'folder') in points_dict:
                          points_dict[(file[0],'folder')]+=1*file[1]
                  else:
                      for word in file:
                          if word in memory:
                              for file2 in memory[word]:
                                  if len(file2)>1:
                                      if 'int' in str(type(file2[1])):
                                          if (file2[0],get_ext(listoffiles[file[0]])) in points_dict:
                                              points_dict[(file2[0],get_ext(file2[0]))]+=1*file2[1]
                                          if (file2[0],'folder') in points_dict:
                                              points_dict[(file2[0],'folder')]+=1*file2[1]
              else:
                 for word in file:
                    if word in memory:
                        for file2 in memory[word]:
                            if len(file2)>1:
                                if 'int' in str(type(file2[1])):
                                    if (file2[0],get_ext(listoffiles[file2[0]])) in points_dict:
                                        points_dict[(file2[0],get_ext(listoffiles[file2[0]]))]+=1*file2[1]
                                    if (file2[0],'folder') in points_dict:
                                        points_dict[(file2[0],'folder')]+=1*file2[1]
   count=0
   rowdr=0;rowm=0;rowds=0
   if alt_count>0:
      frame_4=frame4
      frame_2=frame2
      frme_2=frme2
      frme_=frme
      frame_3=frame3
      frame_1=frame1
   else:
      frame_4=a.frame4
      frame_2=a.frame2
      frme_2=a.frme2
      frme_=a.frme
      frame_3=a.frame3
      frame_1=a.frame1
   for key,value in sorted(points_dict.items(),key=operator.itemgetter(1),reverse=True):
     if alt_UI:
         ldd=Label(alt_Frame,text=key[1]+" File",bg="burlywood2",anchor=E,width=18)
         ldd.grid(row=rowdr,column=0)
         if key[1]=="exe":
             ldd.config(bg="slate gray",text="EXECUTABLE",fg="white")
         elif key[1]=="pdf":
             ldd.config(bg="firebrick1",text="pdf Document",fg="white")
         elif key[1]=="txt":
             ldd.config(bg="antiquewhite1",text="TEXT Document",fg="black")
         elif key[1]=="folder":
             ldd.config(bg="burlywood2",text="FILE FOLDER",fg="black")
         ld=Button(alt_Frame,height=1,anchor=W,text=os.path.basename(listoffiles[key[0]]).split("."+key[1])[0],width=120,bg="white",highlightcolor="blue",relief=GROOVE,command=partial(opennow,listoffiles[key[0]]))
         ld.grid(row=rowdr, column=1)
         ld2=Button(alt_Frame,cursor='hand2',text="Open File Location",command=partial(opennow,os.path.dirname(listoffiles[key[0]])),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
         ld2.grid(row=rowdr,column=2)
         rowdr+=1
         count+=1
         if count==70:
             break
         continue
     if key[1]=="folder":
         folder_count+=1
         ld=Button(frame_4,height=1,anchor=W,text="{}".format(os.path.basename(listoffiles[key[0]])),width=81,bg="navajo white",highlightcolor="blue",relief=GROOVE,command=partial(opennow,listoffiles[key[0]]))
         ld.grid(row=rowdr, column=0)
         ld.bind("<MouseWheel>", partial(_on_mousewheel,frame_3))
         ld2=Button(frame_4,cursor='hand2',text="Location",command=partial(opennow,os.path.dirname(listoffiles[key[0]])),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
         ld2.grid(row=rowdr,column=1)
         rowdr+=1
     elif "mp3" in key[1].lower()  or "3gp" in key[1].lower() or "mp4" in key[1].lower() or "flv" in key[1].lower() or "mkv" in key[1].lower() or "wav" in key[1] or "avi" in key[1] or "jpg" in key[1].lower() or "jpeg" in key[1] or "gif" in key[1] or "png" in key[1]:
         media_count+=1
         lm=Button(frame_2,anchor=W,text="{}".format(os.path.basename(listoffiles[key[0]])),width=81,bg="white",relief=RIDGE,command=partial(opennow,listoffiles[key[0]]))
         lm.grid(row=rowm, column=0)
         lm2=Button(frame_2,cursor='hand2',text="Location",command=partial(opennow,os.path.dirname(listoffiles[key[0]])),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
         lm2.grid(row=rowm,column=1)
         rowm+=1
         lm.bind("<MouseWheel>", partial(_on_mousewheel,frame_1))
     else:
         doc_count+=1
         l=Button(frme_2,anchor=W,text="{}".format(os.path.basename(listoffiles[key[0]])),width=87,bg="white",relief=RIDGE,command=partial(opennow,listoffiles[key[0]]))
         l.grid(row=rowds,column=0)
         l2=Button(frme_2,cursor='hand2',text="Location",command=partial(opennow,os.path.dirname(listoffiles[key[0]])),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
         l2.grid(row=rowds,column=1)
         rowds+=1
         l.bind("<MouseWheel>", partial(_on_mousewheel,frme_))

     count+=1
     if count==100:
         break
 

def BinarySearchMain(list2,first,last,search):
  global iterations
  mid=int((first+last)/2)
  if(first<=last):
    if(list2[mid].startswith(search)):
       listofwords.append(mid)
       traverse(list2,mid,search)
       iterations+=1
    elif (search>list2[mid]):
        BinarySearchMain(list2,mid+1,last,search)
    elif(search<list2[mid]):
        BinarySearchMain(list2,first,mid-1,search)


def traverse(list2,mid,search):
   global iterations
   i=mid
   while(list2[i+1].startswith(search)):
      listofwords.append(i+1)
      iterations+=1
      i+=1
   i=mid
   while(list2[i-1].startswith(search)):
      listofwords.append(i-1)
      iterations+=1
      i-=1

#alternate search function deleted...it used to be here

def search2(event):
    global frame1,frame2,frme,frme2,frame3,frame4,a,thisdict2
    if alt_count>0:
      frme_2=frme2
      frme_=frme
    else:
      frme_2=a.frme2
      frme_=a.frme
    dict3={}
    txt = a.entryBox.get()
    thislist=[]
    if txt !='':
       for element in txt.split():
           elelist=[element.lower()]
           for _ in elelist:
               thislist.append(_)
       for words in thislist:
           try:
              if words in thisdict2 and words!='':
                  if words in dict3:
                      continue
                  dict3[words]=thisdict2[words]
           except:
               continue
       n=len(thislist)
       e=n
       suplist={}
       count=0
       for word in dict3:
           if len(dict3[word])<100:
              for ent in dict3[word]:
                  if listofdocs[ent] in suplist:
                      suplist[listofdocs[ent]]+=1
                      continue
                  suplist[listofdocs[ent]]=1
       rowdr=0;rowds=0
       if alt_UI:
           clearf(alt_Frame)
       else:
           clearf(frme_2)
       for key,value in sorted(suplist.items(),key=operator.itemgetter(1),reverse=True):
           if alt_UI:
               ld=Button(alt_Frame,height=1,anchor=W,text=os.path.basename(key),width=120,bg="white",highlightcolor="blue",relief=GROOVE,command=partial(opennow,os.path.normpath(key)))
               ld.grid(row=rowdr, column=1)
               ld2=Button(alt_Frame,cursor='hand2',text="Open File Location",command=partial(opennow,os.path.dirname(key)),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
               ld2.grid(row=rowdr,column=2)
               ld2.bind("<MouseWheel>", partial(_on_mousewheel,frme_))
               rowdr+=1
           else:
               l=Button(frme_2,anchor=W,text="{}".format(os.path.basename(key)),width=87,bg="white",relief=RIDGE,command=partial(opennow,os.path.normpath(key)))
               l.grid(row=rowds,column=0)
               l2=Button(frme_2,cursor='hand2',text="Location",command=partial(opennow,os.path.dirname(key)),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
               l2.grid(row=rowds,column=1)
               rowds+=1
               l.bind("<MouseWheel>", partial(_on_mousewheel,frme_))


#SEARCH FUNCTIONS END HERE______________________________________________________________________

#===============================================================================

scheck=None
def initSpellCheck():
    global scheck,x
    scheck=spellCheck2.spellCheck(x)


def spellChecking(word):
    global scheck
    print(scheck.correction(word))
    return (scheck.correction(word))

def opennow(path):
    memorize(path)
    os.popen('explorer.exe "{}"'.format(path))

def openfolder(path):
    os.popen('explorer.exe "{}"'.format(os.path.dirname(path)))

def onFrameConfigure(canvas):
    '''Reset the scroll region to encompass the inner frame'''
    canvas.configure(scrollregion=canvas.bbox("all"))

def alt_ui():
   global alt_Frame,alt_Canvas,alt_count
   global frame1,frame2,frme,frme2,frame3,frame4,lm,ld,lf
   if alt_count>0:
      frame_4=frame4
      frame_2=frame2
      frme_2=frme2
      frme_=frme
      frame_3=frame3
      frame_1=frame1
      lm_=lm
      ld_=ld
      lf_=lf
   else:
      frame_4=a.frame4
      frame_2=a.frame2
      frme_2=a.frme2
      frme_=a.frme
      frame_3=a.frame3
      frame_1=a.frame1
      lm_=a.lm
      ld_=a.ld
      lf_=a.lf
   alt_count+=1
   frame_1.destroy()
   frame_3.destroy()
   frme_.destroy()
   lm_.destroy()
   ld_.destroy()
   lf_.destroy()
   alt_Canvas=Canvas(a.root,height=700,width=1360,bg="white",scrollregion=(0,0,1361,701))
   alt_Canvas.pack(side=LEFT)
   alt_Canvas.pack_propagate(False)
   
   alt_Frame=Frame(alt_Canvas,height=700,width=1360,bg="white")
   alt_Frame.pack(side=LEFT)
   vbar=Scrollbar(alt_Canvas,orient=VERTICAL,relief=FLAT,width=18)
   vbar.pack(side=RIGHT,fill=Y)
   vbar.config(command=alt_Canvas.yview)
   alt_Canvas.config( yscrollcommand=vbar.set)
   alt_Canvas.pack(side=LEFT,expand=True,fill=BOTH)
   alt_Canvas.create_window((4,4), window=alt_Frame, anchor="nw")
   alt_Frame.update()
   alt_Frame.bind("<Configure>", lambda event, canvas=alt_Canvas: onFrameConfigure(canvas))
   alt_Canvas.place(x=0,y=70)
   a.root.bind_all("<MouseWheel>", partial(_on_mousewheel,alt_Canvas))


def regular_ui():
        global frame1,frame2,frme,frme2,frame3,frame4,lm,ld,lf
        alt_Canvas.destroy()
        frame1=Canvas(a.root,height=320,width=653,scrollregion=(0,0,350,354))
        frame2=Frame(frame1,height=320,width=653)
        frme=Canvas(a.root,height=640,width=695,scrollregion=(0,0,699,642))
        frme2=Frame(frme,height=640,width=695)
        frame3=Canvas(a.root,height=293,width=653,scrollregion=(0,0,350,351))
        frame4=Frame(frame3,height=293,width=653)
        lm=Label(a.root,text='MEDIA FILES',width=93,anchor=W,bg='white',font='Mincho 10')
        ld=Label(a.root,text='DOCUMENTS',width=99,anchor=W,bg='light gray',font='Mincho 10')
        lf=Label(a.root,text='FOLDERS',width=93,anchor=W,bg='white smoke',font='Mincho 10')
        frame1.config(background='white')
        frame1.pack_propagate(False)
        frame2.pack()
        frame2.configure(background="white")
        vbar=Scrollbar(frame1,orient=VERTICAL,width=14)
        vbar.pack(side=RIGHT,fill=Y)
        vbar.config(command=frame1.yview)
        frame1.config( yscrollcommand=vbar.set)
        frame1.pack(side=LEFT,expand=True,fill=BOTH)
        frame1.create_window((4,4), window=frame2, anchor="nw")
        frame2.update()
        frame2.bind("<Configure>", lambda event, canvas=frame1: onFrameConfigure(canvas))
        frame1.place(x=704,y=85)
        frme.pack()
        frme.config(background='light gray')
        frme2.config(background='light gray')
        frme.pack_propagate(False)
        vbar2=Scrollbar(frme,orient=VERTICAL,relief=FLAT,width=14)
        vbar2.pack(side=RIGHT,fill=Y)
        vbar2.config(command=frme.yview)
        frme.config( yscrollcommand=vbar2.set)
        frme.pack(side=LEFT,expand=True,fill=BOTH)
        frme.create_window((4,4), window=frme2, anchor="nw")
        frme2.update()
        frme2.bind("<Configure>", lambda event, canvas=frme: onFrameConfigure(canvas))
        frme.place(x=3,y=85)
        frame3.pack()
        frame3.pack_propagate(False)
        vbar3=Scrollbar(frame3,orient=VERTICAL,relief=FLAT,width=14)
        vbar3.pack(side=RIGHT,fill=Y)
        vbar3.config(command=frame3.yview)
        frame3.config( yscrollcommand=vbar3.set)
        frame3.pack(side=LEFT,expand=True,fill=BOTH)
        frame3.config(background="white smoke")
        frame4.config(background="white smoke")
        frame3.create_window((4,4), window=frame4, anchor="nw")
        frame4.update()
        frame4.bind("<Configure>", lambda event, canvas=frame3: onFrameConfigure(canvas))
        frame3.place(x=704,y=417)
        ld.pack(ipadx=10)
        ld.place(x=3,y=65)
        lm.pack(ipadx=10)
        lm.place(x=704,y=65)
        lf.pack(ipadx=10)
        lf.place(x=704,y=396)


def change_ui():
   global alt_UI
   alt_UI=not alt_UI
   if alt_UI:
      alt_ui()
   if not alt_UI:
      regular_ui()


def memorize(file):
    global current_query,memory
    current_query_list=current_query.split()
    for _ in current_query.split():
        if _ in memory:
            found=False
            for entry in memory[_]:
                if listoffiles.index(file) in entry:
                    entry[1]+=1
                    found=True
                    break
            if not found:
                memory[_].append([listoffiles.index(file),1])
        else:
            memory[_]=[[listoffiles.index(file),1]]
    for _ in remove(os.path.basename(file).lower()):
       if _ in current_query_list:
          continue
       if _ in memory:
            found=False
            for entry in memory[_]:
                if listoffiles.index(file) in entry:
                    entry[1]+=1
                    found=True
                    break
            if not found:
               if not current_query_list in memory[_]:
                   memory[_].append(current_query_list)
       else:
               memory[_]=[current_query_list]
    memoryfile=open(r'C:\database2\dbm.pkl','wb')
    pickle.dump(memory,memoryfile)
    memoryfile.close()


def load_memory():
    global memory
    try:
       memoryfile=open(r'C:\database2\dbm.pkl','rb')
       memory=pickle.load(memoryfile)
       memoryfile.close()
    except:
        pass


def clearf(frame):
    a=frame.winfo_children()
    for widget in a:
        widget.destroy()


def search():
   global listofwords,points_dict,list2,iterations,a,current_query,shortcut_active,shrt
   lbl=None
   k="";k2=""
   while(True):
     if k!="":
         k2=k
     time.sleep(0.15)
     if shortcut_active:
         eBox=shrt.entrybox
     else:
         eBox=a.entryBox
     try:
       k=str(eBox.get())
       if k2!=k and k2+" "!=k:
         if lbl:
           lbl.destroy()
         listofwords.clear()
         points_dict.clear()
         iterations=0
         current_query=k.lower()
         for _ in k.lower().split():
            iterations=0
            BinarySearchMain(list2,0,len(list2),_)
         dist_qrys=k.lower().split()
         if iterations==0:
            if k!="":
              temp_text=""
              for _ in k.lower().split():
                 temp_text+=spellChecking(_)+" "
                 BinarySearchMain(list2,0,len(list2),spellChecking(_))
              lbl=Label(a.root,text="showing results for {}".format(temp_text),fg="white smoke",width=93,anchor=W,bg='gray53',font='Mincho 14')
              lbl.pack(side=LEFT)
              lbl.place(x=2,y=40)
         process_indx()
       else:
           continue
     except:
       break

def checkInstance(before_shortcut=False):
   try:
      tlist=os.popen('tasklist /FI "IMAGENAME eq searchBASE.exe"')
      mypid2=os.getpid()
      mypid=str(mypid2)
      i=0
      alist=[]
      for line in tlist:
          alist.append(line)
      tlist=os.popen('tasklist /FI "IMAGENAME eq searchBASE.exe"')
      if len(alist)>4:
         for line in tlist:
             i+=1
             if i>3:
                b=line.split()
                print(b[1],'  ',mypid)
                if b[1]==mypid:
                   if not before_shortcut:
                      os.popen("taskkill /F /PID {}".format(b[1]))
                   import updateDbase
                   updateDbase.sentinel=True                   
   except:
       print('No')

def initiate_shortcut():
    global shrt
    thisdict2.clear()
    checkInstance(before_shortcut=True)
    shrt=shortcut.shortcut()
    thrd=threading.Thread(target=searchShort)
    thrd.start()
    shrt.mainloop()


def process_indx2():
   global memory,current_query,shrt
   for _ in listofwords:
      for ent in x[list2[_]]:
         points_dict[ent]=0
         temp_file=os.path.basename(listoffiles[ent[0]]).lower()
         for _2 in current_query.split():
            if _2 in temp_file:
               points_dict[ent]+=1
   for _ in current_query.split():
       if _ in memory:
           for file in memory[_]:
              if len(file)>1:
                if 'int' in str(type(file[1])):
                 if (file[0],get_ext(listoffiles[file[0]])) in points_dict:
                     points_dict[(file[0],get_ext(listoffiles[file[0]]))]+=1.5*file[1]
                 if (file[0],'folder') in points_dict:
                     points_dict[(file[0],'folder')]+=1.5*file[1]
                else:
                 for word in file:
                    if word in memory:
                        for file2 in memory[word]:
                            if len(file2)>1:
                              if 'int' in str(type(file2[1])):
                               if (file2[0],get_ext(listoffiles[file2[0]])) in points_dict:
                                   points_dict[(file2[0],get_ext(listoffiles[file2[0]]))]+=1.5*file2[1]
                               if (file2[0],'folder') in points_dict:
                                   points_dict[(file2[0],'folder')]+=1.5*file2[1]
              else:
                 for word in file:
                    if word in memory:
                        for file2 in memory[word]:
                            if len(file2)>1:
                              if 'int' in str(type(file2[1])):
                               if (file2[0],get_ext(listoffiles[file2[0]])) in points_dict:
                                   points_dict[(file2[0],get_ext(listoffiles[file2[0]]))]+=1.5*file2[1]
                               if (file2[0],'folder') in points_dict:
                                   points_dict[(file2[0],'folder')]+=1.5*file2[1]
   count=0
   rowdr=0;
   mainFrame=shrt.frame
   for key,value in sorted(points_dict.items(),key=operator.itemgetter(1),reverse=True):
         ldd=Label(mainFrame,text=key[1]+" File",bg="burlywood2",anchor=E,width=13)
         ldd.grid(row=rowdr,column=0)
         if key[1]=="exe":
             ldd.config(bg="slate gray",text="EXECUTABLE",fg="white")
         elif key[1]=="pdf":
             ldd.config(bg="firebrick1",text="pdf Document",fg="white")
         elif key[1]=="txt":
             ldd.config(bg="antiquewhite1",text="TEXT Document",fg="black")
         elif key[1]=="folder":
             ldd.config(bg="burlywood2",text="FILE FOLDER",fg="black")
         ld=Button(mainFrame,height=1,anchor=W,text=os.path.basename(listoffiles[key[0]]).split("."+key[1])[0],width=90,bg="white",highlightcolor="blue",relief=GROOVE,command=partial(opennow,listoffiles[key[0]]))
         ld.grid(row=rowdr, column=1)
         ld2=Button(mainFrame,cursor='hand2',text="Open File Location",command=partial(opennow,os.path.dirname(listoffiles[key[0]])),bg="SteelBlue2",relief=FLAT,activebackground="gray53")
         ld2.grid(row=rowdr,column=2)
         rowdr+=1
         count+=1
         if count==20:
             rowdr=0
             count=0
             break
         continue


def searchShort():
   global listofwords,points_dict,list2,iterations,a,current_query,shortcut_active,shrt
   lbl=None
   k="";k2=""
   frame_is_empty=True
   while(True):
     if k!="":
         k2=k
     time.sleep(0.3)
     try:
         if shrt.entrybox.focus_get()==shrt.entrybox:
             k=str(shrt.entrybox.get())
             if k2!=k and k2+" "!=k:
                 if lbl:
                     lbl.destroy()
                 listofwords.clear()
                 points_dict.clear()
                 iterations=0
                 current_query=k.lower()
                 for _ in k.lower().split():
                     iterations=0
                     BinarySearchMain(list2,0,len(list2),_)
                 if iterations==0:
                     if k!="":
                         temp_text=""
                         for _ in k.lower().split():
                             temp_text+=spellChecking(_)+" "
                             BinarySearchMain(list2,0,len(list2),spellChecking(_))
                             shrt.entrybox.config(textvariable=StringVar(value=temp_text))
                 process_indx2()
                 frame_is_empty=False
             else:
                 continue
         else:
             if not frame_is_empty:
                 shrt.frame.destroy()
                 shrt.makeFrame()
                 frame_is_empty=True
     except:
        break


setting=None
def getSettings():
   global setting
   setting=settings.settings()
   setting.mainloop()
   setting=None


def main():
   global listofwords,points_dict,list2,iterations,a,x,shrt,metaclass,supremedir
   supremedir=os.getcwd()
   a=searchbasenew.searchbase()
   if alt_UI:
      alt_ui()
   alt_ui_check=Checkbutton(a.root,command=change_ui,cursor='hand2',text='Alternate UI',bg='gray53' ,fg='black',font='Times 10')
   alt_ui_check.pack()
   alt_ui_check.place(relx=1,x=-80,y=2,anchor=NE)
   os.chdir(supremedir)
   try:
      a.root.iconbitmap('icons\search.ico')
   except:
      try:
         a.root.iconbitmap(r'C:\Program Files\searchBASE\icons\search.ico')#for the icon
      except:
         print("no")
   myButton = Button(a.root, cursor='hand2', text="Search inside Documents", bg="SteelBlue1", relief=GROOVE, activebackground="white", width=19)
   myButton.pack(side=TOP)
   myButton.bind('<Button-1>',search2)
   mymenu=Menu(a.root)
   utils=Menu(mymenu,tearoff=1)
   utils.add_command(label="Settings",activebackground='SteelBlue1',background='peach puff',command=getSettings)
   mymenu.add_cascade(label="Options",menu=utils,activebackground='SteelBlue1',background='peach puff')
   a.root.config(menu=mymenu)
   retrieve_dbase()
   initSpellCheck()
   load_memory()
   try:
      threading.Thread(target=search).start()
   except:
      print("no")
   initgetmeta()
   threading.Thread(target=partial(updateDbase.start,x,most_recent,list2,listofdocs,listoffiles,metaclass)).start()
   a.root.mainloop()


main()


try:
    fileopen2=open(r'C:\database2\DATABASES.pkl','rb')
    x2=pickle.load(fileopen2)
    var=x2["variable"]
    fileopen2.close()
except:
    var=True
if var:
   initiate_shortcut()




checkInstance()
